# -*- coding: utf-8 -*-
"""Shared design system for the Finsight submission deck."""
from pptx import Presentation
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor as C
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE as DASH
from pptx.oxml.ns import qn

# ---------------------------------------------------------------- palette
GROUND      = C(0xF6,0xF7,0xF6)
SURFACE     = C(0xFF,0xFF,0xFF)
SURFACE2    = C(0xF0,0xF2,0xF1)
INK         = C(0x16,0x20,0x2B)
INK2        = C(0x4A,0x57,0x63)
INK3        = C(0x78,0x84,0x8E)
LINE        = C(0xDF,0xE4,0xE3)
LINE2       = C(0xEA,0xED,0xEC)
ACCENT      = C(0x0F,0x6D,0x63)
ACCENT_SOFT = C(0xE3,0xF0,0xEE)
ACCENT_INK  = C(0x0A,0x4F,0x48)
ACCENT_LT   = C(0x4F,0xBC,0xAE)
WARN        = C(0xA8,0x6A,0x10)
WARN_SOFT   = C(0xFB,0xF0,0xDC)
DARK        = C(0x0E,0x17,0x1A)
DARK2       = C(0x18,0x24,0x28)
DARK_INK    = C(0xE6,0xEC,0xEA)
DARK_INK2   = C(0x9F,0xB0,0xAE)

SANS = "Helvetica Neue"
MONO = "Menlo"

W, H = 13.333, 7.5
ML, MR = 0.72, 0.72
CW = W - ML - MR            # 11.893 content width

# ---------------------------------------------------------------- styles
STYLES = {
    "h1":      (30,  INK,       True,  SANS),
    "h1d":     (54,  DARK_INK,  True,  SANS),
    "eyebrow": (10.5,ACCENT,    True,  SANS),
    "sub":     (13,  INK2,      False, SANS),
    "lede":    (12.5,INK2,      False, SANS),
    "body":    (11,  INK2,      False, SANS),
    "bodyi":   (11,  INK,       False, SANS),
    "bold":    (11,  INK,       True,  SANS),
    "accent":  (11,  ACCENT_INK,True,  SANS),
    "warn":    (11,  WARN,      False, SANS),
    "warnb":   (11,  WARN,      True,  SANS),
    "mono":    (9.5, ACCENT_INK,False, MONO),
    "monod":   (9.5, INK2,      False, MONO),
    "muted":   (9.5, INK3,      False, SANS),
    "tiny":    (8.5, INK3,      False, SANS),
    "th":      (10,  INK2,      True,  SANS),
    "big":     (26,  INK,       True,  SANS),
    "bigA":    (26,  ACCENT_INK,True,  SANS),
    "lbl":     (9.5, INK3,      False, SANS),
    "cardh":   (12,  INK,       True,  SANS),
    "q":       (12.5,ACCENT_INK,True,  SANS),
    "a":       (12,  INK,       False, SANS),
}

def _apply(run, style, **ov):
    size, color, bold, font = STYLES[style]
    f = run.font
    f.size = Pt(ov.get("size", size))
    f.color.rgb = ov.get("color", color)
    f.bold = ov.get("bold", bold)
    f.name = ov.get("font", font)
    f.italic = ov.get("italic", False)

# ---------------------------------------------------------------- shapes
def _noshadow(sh):
    try:
        sh.shadow.inherit = False
    except Exception:
        pass

def rect(slide, x, y, w, h, fill=None, line=None, lw=1.0, dash=None,
         radius=None, rounded=True):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shp, In(x), In(y), In(w), In(h))
    if rounded:
        try:
            sh.adjustments[0] = radius if radius is not None else min(0.18, 0.09/max(h,0.01))
        except Exception:
            pass
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
        if dash:
            sh.line.dash_style = dash
    _noshadow(sh)
    sh.text_frame.word_wrap = True
    return sh

def hair(slide, x, y, w, color=LINE, h=0.012):
    return rect(slide, x, y, w, h, fill=color, rounded=False)

def vhair(slide, x, y, h, color=LINE, w=0.012):
    return rect(slide, x, y, w, h, fill=color, rounded=False)

def arrow(slide, x1, y1, x2, y2, color=INK3, lw=1.1, dash=None):
    from pptx.enum.shapes import MSO_CONNECTOR
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, In(x1), In(y1), In(x2), In(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(lw)
    if dash:
        cn.line.dash_style = dash
    ln = cn.line._get_or_add_ln()
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'sm', 'len': 'med'})
    ln.append(tail)
    return cn

# ---------------------------------------------------------------- text
def tf_of(sh, pad=0.10, anchor=MSO_ANCHOR.TOP):
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = In(pad); tf.margin_right = In(pad)
    tf.margin_top = In(pad*0.6); tf.margin_bottom = In(pad*0.6)
    tf.vertical_anchor = anchor
    return tf

def box(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, pad=0.0):
    sh = slide.shapes.add_textbox(In(x), In(y), In(w), In(h))
    return tf_of(sh, pad=pad, anchor=anchor)

def write(tf, paras, align=PP_ALIGN.LEFT):
    """paras: list of (runs, opts). runs = [(text, style, ov?), ...]"""
    first = True
    for item in paras:
        runs, opts = item if isinstance(item, tuple) else (item, {})
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = opts.get("align", align)
        if "space_before" in opts: p.space_before = Pt(opts["space_before"])
        p.space_after = Pt(opts.get("space_after", 0))
        if "ls" in opts: p.line_spacing = opts["ls"]
        if isinstance(runs, str):
            runs = [(runs, opts.get("style", "body"))]
        for r in runs:
            txt, st = r[0], r[1]
            ov = r[2] if len(r) > 2 else {}
            run = p.add_run(); run.text = txt
            _apply(run, st, **ov)
    return tf

def T(slide, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, pad=0.0):
    return write(box(slide, x, y, w, h, anchor=anchor, pad=pad), paras, align)

# ---------------------------------------------------------------- deck
def new_deck():
    prs = Presentation()
    prs.slide_width = In(W); prs.slide_height = In(H)
    return prs

def blank(prs, bg=GROUND):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r = rect(s, -0.05, -0.05, W+0.1, H+0.1, fill=bg, rounded=False)
    return s

def header(slide, eyebrow, title, sub=None, rule=True):
    T(slide, ML, 0.46, CW, 0.28, [([(eyebrow.upper(), "eyebrow", {"size": 10.5})], {})])
    T(slide, ML, 0.72, CW, 0.62, [([(title, "h1")], {})])
    y = 1.40
    if sub:
        T(slide, ML, 1.33, CW*0.86, 0.50, [([(sub, "lede", {"size": 11.5})], {"ls": 1.15})])
        y = 1.82
    if rule:
        hair(slide, ML, y, CW)
    return y + 0.30

def footer(slide, n):
    T(slide, ML, 6.95, 6.0, 0.3, [([("Finsight  ·  TBX / BVP Tech Catalyst", "tiny")], {})])
    T(slide, W-MR-1.4, 6.95, 1.4, 0.3, [([(str(n), "tiny")], {})], align=PP_ALIGN.RIGHT)

# ---------------------------------------------------------------- grid
def grid(slide, x, y, w, col_w, head, rows, row_h, head_h=0.40,
         head_fill=SURFACE2, body_fill=SURFACE, pad=0.16, valign=MSO_ANCHOR.TOP):
    heights = row_h if isinstance(row_h, (list, tuple)) else [row_h]*len(rows)
    total = head_h + sum(heights)
    rect(slide, x, y, w, total, fill=body_fill, line=LINE, lw=0.9, radius=0.035)
    rect(slide, x, y, w, head_h, fill=head_fill, line=None, radius=0.09)
    rect(slide, x, y+head_h-0.10, w, 0.10, fill=head_fill, rounded=False)
    hair(slide, x, y+head_h, w, LINE)
    # header text
    cx = x
    for i, htxt in enumerate(head):
        T(slide, cx+pad, y+0.085, col_w[i]-2*pad, head_h-0.12,
          [([(htxt, "th")], {})])
        cx += col_w[i]
    # rows
    ry = y + head_h
    for ri, row in enumerate(rows):
        cx = x
        for ci, cell in enumerate(row):
            if cell is None:
                cx += col_w[ci]; continue
            paras = cell if isinstance(cell, list) and cell and isinstance(cell[0], (tuple, list)) and isinstance(cell[0][0], (list, tuple)) else None
            if isinstance(cell, str):
                content = [([(cell, "body")], {})]
            elif paras is not None:
                content = cell
            else:
                content = [(cell, {})]
            T(slide, cx+pad, ry+0.085, col_w[ci]-2*pad, heights[ri]-0.12,
              content, anchor=valign)
            cx += col_w[ci]
        ry += heights[ri]
        if ri < len(rows)-1:
            hair(slide, x+0.0, ry, w, LINE2)
    return y + total

# ---------------------------------------------------------------- cards
def statcard(slide, x, y, w, h, big, label, accent=False, fill=SURFACE):
    rect(slide, x, y, w, h, fill=ACCENT_SOFT if accent else fill,
         line=ACCENT if accent else LINE, lw=1.3 if accent else 0.9, radius=0.10)
    T(slide, x+0.18, y+0.16, w-0.36, 0.5,
      [([(big, "bigA" if accent else "big")], {})])
    T(slide, x+0.18, y+0.16+0.46, w-0.36, h-0.62,
      [([(label, "lbl")], {"ls": 1.15})])

def callout(slide, x, y, w, h, runs, fill=WARN_SOFT, bar=WARN, anchor=MSO_ANCHOR.MIDDLE):
    rect(slide, x, y, w, h, fill=fill, line=None, radius=0.06)
    rect(slide, x, y, 0.045, h, fill=bar, rounded=False)
    T(slide, x+0.26, y+0.10, w-0.5, h-0.20, [(runs, {"ls": 1.22})], anchor=anchor)

def chip(slide, x, y, text, fill=ACCENT_SOFT, color=ACCENT_INK, size=8.5, h=0.26, padx=0.14):
    w = 0.085*size/8.5*len(text) + 2*padx
    rect(slide, x, y, w, h, fill=fill, line=None, radius=0.5)
    T(slide, x, y+0.02, w, h, [([(text, "tiny", {"color": color, "size": size, "bold": True})], {})],
      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return x + w + 0.10
