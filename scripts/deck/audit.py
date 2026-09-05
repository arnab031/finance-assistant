# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Emu

import os
P = os.path.normpath(os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "..", "Finsight_Deck.pptx"))
prs = Presentation(P)
SW, SH = prs.slide_width/914400, prs.slide_height/914400

def est(tf, w_in):
    """Estimate rendered text height in inches for a text frame of width w_in."""
    ml = (tf.margin_left or 0)/914400; mr = (tf.margin_right or 0)/914400
    mt = (tf.margin_top or 0)/914400; mb = (tf.margin_bottom or 0)/914400
    avail_pt = max((w_in - ml - mr), 0.05) * 72
    total = 0.0
    for p in tf.paragraphs:
        runs = p.runs
        if not runs:
            total += 8/72.0; continue
        size = max((r.font.size.pt if r.font.size else 12) for r in runs)
        mono = any((r.font.name or "").startswith("Menlo") for r in runs)
        cw = size * (0.60 if mono else 0.505)
        cpl = max(int(avail_pt/cw), 1)
        txt = "".join(r.text for r in runs)
        nlines = max(1, -(-len(txt)//cpl))
        ls = p.line_spacing if isinstance(p.line_spacing, float) else 1.0
        sa = p.space_after.pt if p.space_after else 0
        sb = p.space_before.pt if p.space_before else 0
        total += (nlines * size * 1.21 * ls + sa + sb)/72.0
    return total + mt + mb

bad = []
oob = []
for i, sl in enumerate(prs.slides, 1):
    for sh in sl.shapes:
        x, y = sh.left/914400, sh.top/914400
        w, h = sh.width/914400, sh.height/914400
        if x < -0.06 or y < -0.06 or x+w > SW+0.06 or y+h > SH+0.06:
            oob.append((i, sh.shape_type, round(x,2), round(y,2), round(w,2), round(h,2),
                        (sh.text_frame.text[:40] if sh.has_text_frame else "")))
        if sh.has_text_frame and sh.text_frame.text.strip():
            e = est(sh.text_frame, w)
            if e > h + 0.03:
                bad.append((i, round(e/h,2), round(h,2), round(e,2), round(x,2), round(y,2), round(w,2),
                            sh.text_frame.text[:70].replace("\n"," | ")))
print("== out of bounds ==")
for r in oob: print(r)
print("\n== likely text overflow (est/box) ==")
for r in sorted(bad, key=lambda r:-r[1]): print(r)
print("\nslides:", len(prs.slides._sldIdLst))

print("\n== shapes reaching into the footer band (y+h > 6.90) ==")
for i, sl in enumerate(prs.slides, 1):
    for sh in sl.shapes:
        y = sh.top/914400; h = sh.height/914400; x = sh.left/914400
        if h > 7.0: continue          # background
        txt = sh.text_frame.text[:45].replace("\n"," | ") if sh.has_text_frame else ""
        if txt.startswith("Finsight  ·") or (txt.isdigit() and len(txt)<3): continue
        est_h = est(sh.text_frame, sh.width/914400) if (sh.has_text_frame and txt.strip()) else h
        bot = y + max(h, est_h)
        if bot > 6.90:
            print((i, round(y,2), round(bot,2), txt))
